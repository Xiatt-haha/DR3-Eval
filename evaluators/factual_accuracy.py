#!/usr/bin/env python3
# Copyright 2025 Miromind.ai
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#    http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""

Usage:
    python -m evaluators.factual_accuracy_agent \
        --result examples/001/final_report.md \
        --long-context examples/001/long_context.json \
        --source-folder examples/001/source \
        --output examples/001/eval_factual_accuracy.json
"""

import re
import argparse
import base64
import json
import os
import time
import threading
import subprocess
import tempfile
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed

from dotenv import load_dotenv
from openai import OpenAI

from .utils.base import BaseEvaluator, EvalConfig, EvalResult
from .utils.document_loader import DocumentLoader, normalize_title, fuzzy_title_match
from .utils.llm_client import extract_json_from_text

load_dotenv()

import logging
logger = logging.getLogger("batch_eval")

# API Configuration
API_KEY = os.environ.get("OPENAI_API_KEY", "")
API_BASE_URL = os.environ.get("OPENAI_BASE_URL", "")
MODEL_NAME = os.environ.get("FACTUAL_ACCURACY_MODEL", "gpt-51-1113-global")
GEMINI_MULTIMODAL_MODEL = os.environ.get("GEMINI_MULTIMODAL_MODEL", "gemini-2.5-pro-06-17")
GPT5_MODEL = os.environ.get("GPT5_MODEL", "gpt-51-1113-global")

# File type extensions
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
VIDEO_EXTENSIONS = {".mp4", ".webm", ".mov", ".avi"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac", ".wma"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".html", ".htm", ".md"}
PPTX_EXTENSIONS = {".pptx", ".ppt"}
DATA_EXTENSIONS = {".csv", ".xlsx", ".xls", ".json", ".jsonl", ".tsv"}

# Maximum content length for API calls (in characters, ~4 chars per token)
MAX_CONTENT_CHARS = 20000  # ~5k tokens, reduced to save tokens
BM25_TOP_K = 5
BM25_CHUNK_SIZE = 500

# Rate limiting configuration
API_RPM = int(os.environ.get("FACTUAL_API_RPM", "60"))
API_MAX_RETRIES = int(os.environ.get("FACTUAL_MAX_RETRIES", "3"))
API_RETRY_DELAY = float(os.environ.get("FACTUAL_RETRY_DELAY", "10"))
BATCH_SIZE = 20
MAX_CONCURRENT_REQUESTS = int(os.environ.get("FACTUAL_MAX_CONCURRENT", "1"))

MODEL_MAX_TOKENS_LIMITS = {
    "qwen-max": 8192,
    "qwen-max-latest": 8192,
    "qwen-plus": 8192,
    "qwen-turbo": 8192,
    "glm-4": 4096,
    "glm-4-plus": 4096,
}
DEFAULT_MAX_TOKENS = 16384

def get_max_tokens_for_model(model_name: str) -> int:
    model_name_lower = model_name.lower()
    for model_prefix, limit in MODEL_MAX_TOKENS_LIMITS.items():
        if model_prefix in model_name_lower:
            return limit
    return DEFAULT_MAX_TOKENS

def encode_file_to_base64(file_path: str) -> str:
    with open(file_path, "rb") as f:
        return base64.standard_b64encode(f.read()).decode("utf-8")

def parse_citation_details(citation: str) -> Dict:
    """
    
    - [Doc: llama.pdf, Page 34]
    - [vision_understanding_advanced: image2.jpg, reasoning]
    - [RAG-1]
    
    Returns:
        Dict with keys: source_type, source_name, page, chunk, extra_info
    """
    result = {
        'source_type': 'unknown',
        'source_name': '',
        'page': None,
        'page_range': None,  # (start, end) tuple
        'chunk': None,
        'chunk_range': None,  # (start, end) tuple
        'extra_info': None,
        'raw_citation': citation
    }
    
    content = citation.strip()
    if content.startswith('[') and content.endswith(']'):
        content = content[1:-1]
    elif content.startswith('(') and content.endswith(')'):
        content = content[1:-1]
    
    if 'long_context:' in content:
        result['source_type'] = 'long_context'
        match = re.search(r'long_context:\s*"([^"]+)"', content)
        if match:
            result['source_name'] = match.group(1)
        chunk_match = re.search(r'chunk\s*(\d+)(?:\s*-\s*(\d+))?', content, re.I)
        if chunk_match:
            start = int(chunk_match.group(1))
            end = int(chunk_match.group(2)) if chunk_match.group(2) else start
            result['chunk'] = start
            result['chunk_range'] = (start, end)
    
    elif 'Doc:' in content or '文档:' in content or 'doc:' in content.lower():
        result['source_type'] = 'doc'
        match = re.search(r'(?:[Dd]oc|文档):\s*([^,\]]+)', content)
        if match:
            result['source_name'] = match.group(1).strip()
        page_match = re.search(r'Page\s*(\d+)(?:\s*-\s*(\d+))?', content, re.I)
        if page_match:
            start = int(page_match.group(1))
            end = int(page_match.group(2)) if page_match.group(2) else start
            result['page'] = start
            result['page_range'] = (start, end)
    
    elif 'RAG-' in content:
        result['source_type'] = 'rag'
        match = re.search(r'RAG-(\d+)', content)
        if match:
            result['source_name'] = match.group(1)
    
    elif '图片:' in content or 'image:' in content.lower():
        result['source_type'] = 'image'
        match = re.search(r'(?:图片|[Ii]mage):\s*([^,\]]+)', content)
        if match:
            result['source_name'] = match.group(1).strip()
    
    elif '视频:' in content or 'video:' in content.lower():
        result['source_type'] = 'video'
        match = re.search(r'(?:视频|[Vv]ideo):\s*([^,\]]+)', content)
        if match:
            result['source_name'] = match.group(1).strip()
    
    elif '音频:' in content or 'audio:' in content.lower():
        result['source_type'] = 'audio'
        match = re.search(r'(?:音频|[Aa]udio):\s*([^,\]]+)', content)
        if match:
            result['source_name'] = match.group(1).strip()
    
    elif ':' in content:
        parts = content.split(':', 1)
        result['source_type'] = parts[0].strip().lower()
        if len(parts) > 1:
            remaining = parts[1].strip()
            if ',' in remaining:
                name_parts = remaining.split(',', 1)
                result['source_name'] = name_parts[0].strip()
                result['extra_info'] = name_parts[1].strip()
            else:
                result['source_name'] = remaining
        
        ext_match = re.search(r'\.(\w+)(?:\s|,|$)', result['source_name'])
        if ext_match:
            ext = '.' + ext_match.group(1).lower()
            if ext in IMAGE_EXTENSIONS:
                result['source_type'] = 'image'
            elif ext in VIDEO_EXTENSIONS:
                result['source_type'] = 'video'
            elif ext in AUDIO_EXTENSIONS:
                result['source_type'] = 'audio'
            elif ext in DOCUMENT_EXTENSIONS:
                result['source_type'] = 'doc'
    
    else:
        ext_match = re.search(r'\.(\w+)$', content)
        if ext_match:
            ext = '.' + ext_match.group(1).lower()
            result['source_name'] = content
            if ext in IMAGE_EXTENSIONS:
                result['source_type'] = 'image'
            elif ext in VIDEO_EXTENSIONS:
                result['source_type'] = 'video'
            elif ext in AUDIO_EXTENSIONS:
                result['source_type'] = 'audio'
            elif ext in DOCUMENT_EXTENSIONS:
                result['source_type'] = 'doc'
    
    return result

def get_audio_duration(input_path: str) -> float:
    cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration",
           "-of", "default=noprint_wrappers=1:nokey=1", input_path]
    result = subprocess.run(cmd, capture_output=True, text=True)
    try:
        return float(result.stdout.strip())
    except ValueError:
        return 0.0

def compress_audio(input_path: str, output_path: str, max_size_mb: float = 5.0) -> str:
    duration = get_audio_duration(input_path)
    if duration <= 0:
        target_bitrate = 48
    else:
        target_size_bits = max_size_mb * 1024 * 1024 * 8
        target_bitrate = int(target_size_bits / duration / 1000)
        target_bitrate = max(32, min(target_bitrate, 128))
    cmd = ["ffmpeg", "-y", "-i", input_path, "-b:a", f"{target_bitrate}k", "-ar", "22050", output_path]
    subprocess.run(cmd, capture_output=True)
    return output_path

def prepare_audio_for_gemini(filepath: Path) -> Tuple[Optional[List[Dict]], str]:
    ext = filepath.suffix.lower()
    filename = filepath.name
    size_mb = os.path.getsize(filepath) / 1024 / 1024
    content_parts = []
    max_audio_size_mb = 5.0
    audio_path = str(filepath)
    
    if size_mb > max_audio_size_mb:
        logger.info(f"    Compressing audio ({size_mb:.2f} MB)...")
        compressed_path = tempfile.mktemp(suffix=".mp3")
        compress_audio(audio_path, compressed_path, max_size_mb=max_audio_size_mb)
        audio_path = compressed_path
        size_mb = os.path.getsize(audio_path) / 1024 / 1024
        logger.info(f"    Compressed to {size_mb:.2f} MB")
        mime_type = "audio/mp3"
    else:
        mime_types = {".mp3": "audio/mp3", ".wav": "audio/wav", ".m4a": "audio/m4a",
                      ".aac": "audio/aac", ".ogg": "audio/ogg", ".flac": "audio/flac"}
        mime_type = mime_types.get(ext, "audio/mp3")
    
    audio_base64 = encode_file_to_base64(audio_path)
    content_parts.append({"type": "audio_url", "audio_url": {"url": f"data:{mime_type};base64,{audio_base64}"}})
    return content_parts, f"音频文件: {filename} ({size_mb:.2f} MB)"

def compress_video(input_path: str, output_path: str, max_size_mb: float = 5.0) -> str:
    cmd = [
        "ffmpeg", "-y", "-i", input_path,
        "-vf", "scale=320:-2",
        "-c:v", "libx264", "-crf", "32",
        "-preset", "fast",
        "-an",
        output_path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    
    if os.path.exists(output_path):
        output_size_mb = os.path.getsize(output_path) / 1024 / 1024
        logger.info(f"    Compressed video size: {output_size_mb:.2f} MB")
        
        if output_size_mb > max_size_mb:
            logger.warning(f"    Still too large, trying more aggressive compression...")
            temp_output = output_path + ".tmp.mp4"
            cmd = [
                "ffmpeg", "-y", "-i", input_path,
                "-vf", "scale=240:-2",
                "-c:v", "libx264", "-crf", "40",
                "-preset", "fast",
                "-an",
                temp_output
            ]
            subprocess.run(cmd, capture_output=True)
            if os.path.exists(temp_output):
                os.replace(temp_output, output_path)
    
    return output_path

def prepare_video_for_gemini(filepath: Path) -> Tuple[Optional[List[Dict]], str]:
    ext = filepath.suffix.lower()
    filename = filepath.name
    size_mb = os.path.getsize(filepath) / 1024 / 1024
    content_parts = []
    max_video_size_mb = 5.5
    video_path = str(filepath)
    
    if size_mb > max_video_size_mb:
        logger.info(f"    Compressing video ({size_mb:.2f} MB)...")
        compressed_path = tempfile.mktemp(suffix=".mp4")
        try:
            compress_video(video_path, compressed_path, max_size_mb=max_video_size_mb)
            video_path = compressed_path
            size_mb = os.path.getsize(video_path) / 1024 / 1024
            logger.info(f"    Compressed to {size_mb:.2f} MB")
        except Exception as e:
            logger.error(f"    ❌ Video compression failed: {e}")
            return None, f"视频压缩失败: {filename}"
    
    mime_types = {".mp4": "video/mp4", ".webm": "video/webm", ".mov": "video/quicktime", ".avi": "video/x-msvideo"}
    mime_type = mime_types.get(ext, "video/mp4")
    
    video_base64 = encode_file_to_base64(video_path)
    content_parts.append({"type": "video_url", "video_url": {"url": f"data:{mime_type};base64,{video_base64}"}})
    return content_parts, f"视频文件: {filename} ({size_mb:.2f} MB)"

def extract_pptx_pages(filepath: Path, page_numbers: Optional[List[int]] = None) -> str:
    """
    
    Args:
    
    Returns:
    """
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        
        content_parts = []
        total_slides = len(prs.slides)
        
        if page_numbers:
            valid_pages = [p for p in page_numbers if 1 <= p <= total_slides]
            if not valid_pages:
                valid_pages = list(range(1, min(total_slides + 1, 6)))
        else:
            valid_pages = list(range(1, total_slides + 1))
        
        for page_num in valid_pages:
            slide_idx = page_num - 1
            try:
                slide = prs.slides[slide_idx]
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    content_parts.append(f"=== Slide {page_num} ===\n" + "\n".join(slide_text))
            except Exception as e:
                logger.warning(f"  ⚠️ Cannot extract slide {page_num}: {e}")
                continue
        
        return "\n\n".join(content_parts)
    except ImportError:
        logger.error("  ❌ python-pptx not installed, cannot extract PPTX content")
        return ""
    except Exception as e:
        logger.error(f"  ❌ Error extracting PPTX: {e}")
        return ""

def extract_csv_content(filepath: Path, max_rows: int = 100) -> str:
    """
    
    Args:
    
    Returns:
    """
    try:
        import csv
        
        content_parts = []
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= max_rows:
                    content_parts.append(f"... (共 {i}+ 行，已截断)")
                    break
                content_parts.append(" | ".join(row))
        
        content = "\n".join(content_parts)
        
        if len(content) > MAX_CONTENT_CHARS:
            logger.warning(f"  ⚠️ CSV content too long ({len(content)} chars), truncating to {MAX_CONTENT_CHARS}")
            content = content[:MAX_CONTENT_CHARS] + "\n\n[... 内容已截断 ...]"
        
        return content
    except Exception as e:
        logger.error(f"  ❌ Error extracting CSV: {e}")
        return ""

def extract_excel_content(filepath: Path, max_rows: int = 100) -> str:
    """
    
    Args:
    
    Returns:
    """
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        content_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_content = [f"=== Sheet: {sheet_name} ==="]
            
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= max_rows:
                    sheet_content.append(f"... (共 {i}+ 行，已截断)")
                    break
                row_str = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_str.strip():
                    sheet_content.append(row_str)
            
            content_parts.append("\n".join(sheet_content))
        
        wb.close()
        content = "\n\n".join(content_parts)
        
        if len(content) > MAX_CONTENT_CHARS:
            logger.warning(f"  ⚠️ Excel content too long ({len(content)} chars), truncating to {MAX_CONTENT_CHARS}")
            content = content[:MAX_CONTENT_CHARS] + "\n\n[... 内容已截断 ...]"
        
        return content
    except ImportError:
        logger.error("  ❌ openpyxl not installed, cannot extract Excel content")
        return ""
    except Exception as e:
        logger.error(f"  ❌ Error extracting Excel: {e}")
        return ""

def extract_docx_pages(filepath: Path, page_numbers: Optional[List[int]] = None) -> str:
    """
    
    Args:
    
    Returns:
    """
    try:
        from docx import Document
        doc = Document(filepath)
        
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        
        content = "\n\n".join(paragraphs)
        
        if len(content) > MAX_CONTENT_CHARS:
            logger.warning(f"  ⚠️ DOCX content too long ({len(content)} chars), truncating to {MAX_CONTENT_CHARS}")
            content = content[:MAX_CONTENT_CHARS] + "\n\n[... 内容已截断 ...]"
        
        return content
    except ImportError:
        logger.error("  ❌ python-docx not installed, cannot extract DOCX content")
        return ""
    except Exception as e:
        logger.error(f"  ❌ Error extracting DOCX: {e}")
        return ""

def get_page_numbers_from_claims(claims: List[Dict]) -> List[int]:
    """
    
    Args:
    
    Returns:
    """
    page_numbers = set()
    for claim in claims:
        citation = claim.get('citation', '')
        if citation:
            details = parse_citation_details(citation)
            if details.get('page'):
                page_numbers.add(details['page'])
            if details.get('page_range'):
                start, end = details['page_range']
                for p in range(start, end + 1):
                    page_numbers.add(p)
    
    return sorted(list(page_numbers))

def bm25_retrieve_relevant_chunks(content: str, query: str, chunk_size: int = BM25_CHUNK_SIZE, top_k: int = BM25_TOP_K) -> str:
    """
    
    Args:
    
    Returns:
    """
    try:
        from rank_bm25 import BM25Okapi
        import jieba
    except ImportError:
        logger.warning("  ⚠️ rank_bm25 or jieba not installed, falling back to truncation")
        return content[:MAX_CONTENT_CHARS] if len(content) > MAX_CONTENT_CHARS else content
    
    chunks = []
    for i in range(0, len(content), chunk_size):
        chunk = content[i:i + chunk_size]
        if chunk.strip():
            chunks.append(chunk)
    
    if not chunks:
        return content[:MAX_CONTENT_CHARS] if len(content) > MAX_CONTENT_CHARS else content
    
    tokenized_chunks = [list(jieba.cut(chunk)) for chunk in chunks]
    
    bm25 = BM25Okapi(tokenized_chunks)
    
    tokenized_query = list(jieba.cut(query))
    
    scores = bm25.get_scores(tokenized_query)
    
    top_indices = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:top_k]
    
    top_indices = sorted(top_indices)
    
    relevant_chunks = []
    for idx in top_indices:
        chunk_num = idx + 1
        relevant_chunks.append(f"[片段 {chunk_num}]\n{chunks[idx]}")
    
    result = "\n\n".join(relevant_chunks)
    logger.info(f"    🔍 BM25: selected {len(top_indices)} chunks from {len(chunks)} total (query: {query[:30]}...)")
    
    return result

def bm25_retrieve_for_claims(content: str, claims: List[Dict], chunk_size: int = BM25_CHUNK_SIZE, top_k: int = BM25_TOP_K) -> str:
    """
    
    Args:
    
    Returns:
    """
    combined_query = " ".join([c.get('claim_text', '') for c in claims])
    return bm25_retrieve_relevant_chunks(content, combined_query, chunk_size, top_k)

def prepare_file_for_api(filepath: Path) -> Tuple[Optional[List[Dict]], str, str]:
    ext = filepath.suffix.lower()
    filename = filepath.name
    content_parts = []
    api_type = 'chat'
    
    if ext in IMAGE_EXTENSIONS:
        image_base64 = encode_file_to_base64(str(filepath))
        mime_types = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", 
                      ".gif": "image/gif", ".webp": "image/webp"}
        mime_type = mime_types.get(ext, "image/jpeg")
        content_parts.append({"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_base64}"}})
        description = f"图片文件: {filename}"
    elif ext == ".pdf":
        pdf_base64 = encode_file_to_base64(str(filepath))
        content_parts.append({"type": "input_file", "filename": filename, "file_data": f"data:application/pdf;base64,{pdf_base64}"})
        description = f"PDF文档: {filename}"
        api_type = 'responses'
    elif ext in VIDEO_EXTENSIONS:
        size_mb = os.path.getsize(filepath) / 1024 / 1024
        max_video_size_mb = 5.5
        video_path = str(filepath)
        
        if size_mb > max_video_size_mb:
            logger.info(f"    Compressing video ({size_mb:.2f} MB)...")
            compressed_path = tempfile.mktemp(suffix=".mp4")
            try:
                compress_video(video_path, compressed_path, max_size_mb=max_video_size_mb)
                video_path = compressed_path
                size_mb = os.path.getsize(video_path) / 1024 / 1024
                logger.info(f"    Compressed to {size_mb:.2f} MB")
            except Exception as e:
                logger.error(f"    ❌ Video compression failed: {e}")
                return None, f"视频压缩失败: {filename}", api_type
        
        video_base64 = encode_file_to_base64(video_path)
        mime_types = {".mp4": "video/mp4", ".webm": "video/webm", ".mov": "video/quicktime", ".avi": "video/x-msvideo"}
        mime_type = mime_types.get(ext, "video/mp4")
        content_parts.append({"type": "video_url", "video_url": {"url": f"data:{mime_type};base64,{video_base64}"}})
        description = f"视频文件: {filename} ({size_mb:.2f} MB)"
        api_type = 'chat'
    elif ext in {".txt", ".md", ".html", ".htm"}:
        try:
            text_content = filepath.read_text(encoding="utf-8", errors="ignore")
            content_parts.append({"type": "text", "text": f"文档内容 ({filename}):\n\n{text_content}"})
            description = f"文本文档: {filename}"
        except Exception as e:
            description = f"无法读取文档: {filename} ({e})"
    else:
        description = f"不支持的文件类型: {filename}"
    
    return content_parts if content_parts else None, description, api_type

class FactualAccuracyAgentEvaluator(BaseEvaluator):
    metric_name = "factual_accuracy"
    weight = 1.0
    
    def __init__(self, config: Optional[EvalConfig] = None):
        super().__init__(config)
        self.doc_loader = DocumentLoader()
        self.api_client = None
        self.gemini_client = None
        if API_KEY and API_BASE_URL:
            self.api_client = OpenAI(api_key=API_KEY, base_url=API_BASE_URL)
            self.gemini_client = self.api_client
    
    def _load_existing_evaluation(self, case_dir: Path) -> Optional[Dict]:
        """Load existing evaluation results from a case directory."""
        eval_files = list(case_dir.glob("eval_factual_accuracy*.json"))
        if not eval_files:
            return None
        
        eval_files.sort(key=lambda x: x.name, reverse=True)
        eval_file = eval_files[0]
        
        try:
            with open(eval_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
            logger.info(f"  📂 Loaded existing evaluation from: {eval_file.name}")
            return data
        except Exception as e:
            logger.warning(f"  ⚠️ Failed to load existing evaluation: {e}")
            return None
    
    def _find_failed_claims(self, existing_data: Dict) -> Tuple[List[Dict], List[Dict]]:
        """Find claims that failed verification in existing evaluation data."""
        failed_claims = []
        successful_claims = []
        
        if not existing_data:
            return failed_claims, successful_claims
        
        details = existing_data.get('details', {})
        verification_result = details.get('verification_result', {})
        verifications = verification_result.get('verifications', [])
        
        extraction = details.get('extraction', {})
        original_claims = extraction.get('claims', [])
        
        claim_id_to_original = {c.get('id'): c for c in original_claims}
        
        error_patterns = [
            'parse error', 'api error', 'not verified', 'connection error',
            'timeout', 'max retries', 'cannot prepare', 'compression failed',
            'no result', 'empty response'
        ]
        
        for verification in verifications:
            claim_id = verification.get('id')
            explanation = verification.get('explanation', '').lower()
            supported = verification.get('supported', False)
            
            has_error = any(pattern in explanation for pattern in error_patterns)
            
            citation_results = verification.get('citation_results', [])
            for cr in citation_results:
                cr_explanation = cr.get('explanation', '').lower()
                if any(pattern in cr_explanation for pattern in error_patterns):
                    has_error = True
                    break
            
            if has_error:
                original_claim = claim_id_to_original.get(claim_id, {})
                failed_claims.append({
                    'id': claim_id,
                    'claim_text': verification.get('claim_text', original_claim.get('claim_text', '')),
                    'citations': original_claim.get('citations', []),
                    'source_type': original_claim.get('source_type', 'unknown'),
                    'original_verification': verification
                })
            else:
                successful_claims.append(verification)
        
        return failed_claims, successful_claims
    
    def _merge_verification_results(
        self, 
        successful_claims: List[Dict], 
        new_verifications: List[Dict],
        original_claims: List[Dict]
    ) -> Tuple[List[Dict], float]:
        """Merge successful and new verification results."""
        verification_map = {}
        
        for v in successful_claims:
            verification_map[v.get('id')] = v
        
        for v in new_verifications:
            verification_map[v.get('id')] = v
        
        merged = sorted(verification_map.values(), key=lambda x: x.get('id', 0))
        
        supported_count = sum(1 for v in merged if v.get('supported', False))
        total_count = len(merged)
        accuracy_score = (supported_count / total_count * 100) if total_count > 0 else 0
        
        return merged, accuracy_score
    
    def evaluate_batch(self, batch_data: List[Dict]) -> Dict[str, EvalResult]:
        """
        
        Args:
            batch_data: List of dicts, each containing:
                - case_id: str
                - result_text: str
                - long_context_path: Optional[Path]
                - source_folder: Optional[Path]
        
        Returns:
            Dict[case_id, EvalResult]
        """
        import time
        eval_start_time = time.time()
        
        logger.info("=" * 60)
        logger.info(f"🚀 BATCH EVALUATION: {len(batch_data)} reports")
        logger.info("=" * 60)
        
        logger.info("Phase 1: Extracting claims from all reports...")
        phase1_start = time.time()
        
        all_claims_by_case = {}  # {case_id: [claims]}
        for data in batch_data:
            case_id = data['case_id']
            result_text = data['result_text']
            long_context_path = data.get('long_context_path')
            source_folder = data.get('source_folder')
            
            if long_context_path:
                self.doc_loader.load_long_context(long_context_path)
            if source_folder:
                self.doc_loader.load_source_folder(source_folder)
            
            logger.info(f"  [{case_id}] Extracting claims...")
            claims = self._extract_claims_with_llm(result_text)
            
            if not claims:
                logger.warning(f"  [{case_id}] ⚠️ No claims extracted")
                all_claims_by_case[case_id] = []
            else:
                logger.info(f"  [{case_id}] ✅ Extracted {len(claims)} claims")
                all_claims_by_case[case_id] = claims
        
        phase1_time = time.time() - phase1_start
        total_claims = sum(len(claims) for claims in all_claims_by_case.values())
        logger.info(f"Phase 1 completed in {phase1_time:.1f}s")
        logger.info(f"Total claims extracted: {total_claims}")
        
        logger.info("=" * 60)
        logger.info("Phase 2: Global batch verification...")
        phase2_start = time.time()
        
        global_tasks = []  # [(case_id, claim_id, claim_text, citation, source_key)]
        for case_id, claims in all_claims_by_case.items():
            source_folder = next((d['source_folder'] for d in batch_data if d['case_id'] == case_id), None)
            for claim in claims:
                claim_id = claim.get('id')
                claim_text = claim.get('claim_text', '')
                citations = claim.get('citations', [])
                source_type = claim.get('source_type', 'unknown')
                
                if not citations:
                    global_tasks.append((case_id, claim_id, claim_text, None, 'unknown', source_folder))
                else:
                    for citation in citations:
                        source_key = self._get_source_key([citation], source_type)
                        global_tasks.append((case_id, claim_id, claim_text, citation, source_key, source_folder))
        
        logger.info(f"  📋 Created {len(global_tasks)} verification tasks across {len(batch_data)} reports")
        
        source_groups = {}  # {source_key: [(case_id, claim_id, claim_text, citation)]}
        for case_id, claim_id, claim_text, citation, source_key, source_folder in global_tasks:
            if source_key not in source_groups:
                source_groups[source_key] = []
            source_groups[source_key].append((case_id, claim_id, claim_text, citation, source_folder))
        
        logger.info(f"  📚 Grouped into {len(source_groups)} unique source documents")
        
        for source_key, tasks in sorted(source_groups.items(), key=lambda x: len(x[1]), reverse=True)[:10]:
            logger.info(f"    - {source_key[:50]}: {len(tasks)} claims")
        
        task_results = {}  # {(case_id, claim_id, citation): result}
        batch_jobs = []
        
        for source_key, group_tasks in source_groups.items():
            first_source_folder = group_tasks[0][4] if group_tasks else None
            source_info = self._get_source_info(source_key, first_source_folder)
            
            for i in range(0, len(group_tasks), BATCH_SIZE):
                batch = group_tasks[i:i + BATCH_SIZE]
                batch_jobs.append((source_key, source_info, batch))
        
        logger.info(f"  🔄 Created {len(batch_jobs)} batch jobs (max {BATCH_SIZE} claims each)")
        
        def verify_single_batch(job):
            source_key, source_info, batch = job
            temp_claims = [{'id': idx + 1, 'claim_text': t[2], 'citation': t[3]} for idx, t in enumerate(batch)]
            
            thread_client = OpenAI(api_key=API_KEY, base_url=API_BASE_URL)
            
            source_type = source_info.get('type', '')
            file_path = source_info.get('file_path')
            file_ext = file_path.suffix.lower() if file_path else ''
            
            model_supports_image = 'gpt-5' in MODEL_NAME.lower() or 'gpt5' in MODEL_NAME.lower()
            
            if source_type == 'audio':
                results = self._verify_audio_with_gemini_thread(temp_claims, source_info, thread_client)
            elif source_type == 'video':
                results = self._verify_video_with_gemini_thread(temp_claims, source_info, thread_client)
            elif source_type == 'file' and file_ext in IMAGE_EXTENSIONS:
                if model_supports_image:
                    results = self._verify_batch_with_api_thread(temp_claims, source_info, thread_client)
                else:
                    results = self._verify_image_with_gemini_thread(temp_claims, source_info, thread_client)
            elif source_type == 'file' and file_ext == '.pdf':
                results = self._verify_pdf_with_gpt5_thread(temp_claims, source_info, thread_client)
            else:
                results = self._verify_batch_with_api_thread(temp_claims, source_info, thread_client)
            
            batch_results = []
            for idx, (case_id, claim_id, claim_text, citation, _) in enumerate(batch):
                result = results[idx] if idx < len(results) else {'supported': False, 'explanation': 'No result'}
                task_key = (case_id, claim_id, citation or '')
                batch_results.append((task_key, {
                    'supported': result.get('supported', False),
                    'explanation': result.get('explanation', ''),
                    'source_found': result.get('source_found', source_info.get('found', False))
                }))
            return source_key, batch_results
        
        max_workers = min(len(batch_jobs), MAX_CONCURRENT_REQUESTS)
        logger.info(f"  ⚡ Starting parallel verification with {max_workers} workers (max {MAX_CONCURRENT_REQUESTS})")
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(verify_single_batch, job): job for job in batch_jobs}
            batch_num = 0
            for future in as_completed(futures):
                batch_num += 1
                try:
                    source_key, batch_results = future.result()
                    logger.info(f"  Batch {batch_num}/{len(batch_jobs)}: ✅ '{source_key[:40]}...' ({len(batch_results)} claims)")
                    for task_key, result in batch_results:
                        task_results[task_key] = result
                except Exception as e:
                    logger.error(f"  Batch {batch_num} failed: {e}")
        
        phase2_time = time.time() - phase2_start
        logger.info(f"Phase 2 completed in {phase2_time:.1f}s")
        
        logger.info("=" * 60)
        logger.info("Phase 3: Distributing results to reports...")
        
        results_by_case = {}
        for case_id, claims in all_claims_by_case.items():
            if not claims:
                results_by_case[case_id] = EvalResult(
                    metric_name=self.metric_name,
                    score=-1,
                    details={'method': 'llm_agent_extraction', 'error': 'No claims extracted', 'claim_count': 0},
                    weight=self.weight
                )
                continue
            
            verification_results = []
            for claim in claims:
                claim_id = claim.get('id')
                citations = claim.get('citations', [])
                
                if not citations:
                    verification_results.append({
                        'id': claim_id,
                        'claim_text': claim.get('claim_text', '')[:100],
                        'supported': False,
                        'source_found': False,
                        'explanation': 'No citations',
                        'citation_results': []
                    })
                    continue
                
                citation_results = []
                all_supported = True
                any_source_found = False
                explanations = []
                
                for citation in citations:
                    task_key = (case_id, claim_id, citation)
                    result = task_results.get(task_key, {'supported': False, 'explanation': 'Not verified'})
                    citation_results.append({
                        'citation': citation,
                        'supported': result.get('supported', False),
                        'explanation': result.get('explanation', '')
                    })
                    if not result.get('supported', False):
                        all_supported = False
                    if result.get('source_found', False):
                        any_source_found = True
                    explanations.append(f"{citation}: {result.get('explanation', '')[:30]}")
                
                verification_results.append({
                    'id': claim_id,
                    'claim_text': claim.get('claim_text', '')[:100],
                    'supported': all_supported,
                    'source_found': any_source_found,
                    'explanation': '; '.join(explanations),
                    'citation_results': citation_results
                })
            
            supported_count = sum(1 for r in verification_results if r.get('supported'))
            total_claims = len(claims)
            accuracy_score = (supported_count / total_claims * 100) if total_claims > 0 else 0
            
            long_context_claims = []
            source_claims = []
            
            for i, claim in enumerate(claims):
                source_type = claim.get('source_type', 'unknown')
                verification = verification_results[i] if i < len(verification_results) else {}
                
                if source_type in ('long_context', 'rag'):
                    long_context_claims.append({
                        'claim': claim,
                        'verification': verification
                    })
                else:
                    source_claims.append({
                        'claim': claim,
                        'verification': verification
                    })
            
            lc_supported = sum(1 for c in long_context_claims if c['verification'].get('supported'))
            lc_total = len(long_context_claims)
            lc_score = (lc_supported / lc_total * 100) if lc_total > 0 else -1
            
            src_supported = sum(1 for c in source_claims if c['verification'].get('supported'))
            src_total = len(source_claims)
            src_score = (src_supported / src_total * 100) if src_total > 0 else -1
            
            details = {
                'method': 'llm_agent_extraction',
                'extraction': {'total_claims': total_claims, 'claims': claims},
                'verification_result': {
                    'verifications': verification_results,
                    'supported_count': supported_count,
                    'unsupported_count': total_claims - supported_count,
                    'accuracy_score': accuracy_score
                },
                'score_breakdown': {
                    'long_context': {
                        'score': lc_score,
                        'supported_count': lc_supported,
                        'total_count': lc_total,
                        'claims': [c['claim'].get('id') for c in long_context_claims]
                    },
                    'source': {
                        'score': src_score,
                        'supported_count': src_supported,
                        'total_count': src_total,
                        'claims': [c['claim'].get('id') for c in source_claims]
                    }
                }
            }
            
            results_by_case[case_id] = EvalResult(
                metric_name=self.metric_name,
                score=accuracy_score,
                details=details,
                weight=self.weight
            )
            
            lc_score_str = f"{lc_score:.1f}" if lc_score >= 0 else "N/A"
            src_score_str = f"{src_score:.1f}" if src_score >= 0 else "N/A"
            logger.info(f"  [{case_id}] Score: {accuracy_score:.1f}/100 ({supported_count}/{total_claims})")
            logger.info(f"    - Long Context: {lc_score_str}/100 ({lc_supported}/{lc_total})")
            logger.info(f"    - Source Files: {src_score_str}/100 ({src_supported}/{src_total})")
        
        total_time = time.time() - eval_start_time
        logger.info("=" * 60)
        logger.info(f"📊 BATCH EVALUATION COMPLETE")
        logger.info(f"   Reports: {len(batch_data)}")
        logger.info(f"   Total Claims: {total_claims}")
        logger.info(f"   Phase 1 (Extract): {phase1_time:.1f}s")
        logger.info(f"   Phase 2 (Verify):  {phase2_time:.1f}s")
        logger.info(f"   Total Time: {total_time:.1f}s")
        logger.info(f"   Avg per report: {total_time/len(batch_data):.1f}s")
        logger.info("=" * 60)
        
        return results_by_case
    
    def evaluate(self, result_text: str, long_context_path: Optional[Path] = None,
                 source_folder: Optional[Path] = None, case_id: str = None, **kwargs) -> EvalResult:
        """
        
        """
        import time
        eval_start_time = time.time()
        
        if case_id:
            logger.info(f"[{case_id}] Starting evaluation...")
        
        if long_context_path:
            self.doc_loader.load_long_context(long_context_path)
        if source_folder:
            self.doc_loader.load_source_folder(source_folder)
        
        existing_data = None
        if source_folder:
            existing_data = self._load_existing_evaluation(source_folder)
        
        claims = None
        failed_claims = []
        successful_claims = []
        
        if existing_data:
            details = existing_data.get('details', {})
            extraction = details.get('extraction', {})
            original_claims = extraction.get('claims', [])
            
            if original_claims:
                logger.info(f"  📂 Found existing evaluation with {len(original_claims)} claims")
                claims = original_claims
                
                failed_claims, successful_claims = self._find_failed_claims(existing_data)
                
                if failed_claims:
                    logger.info(f"  🔄 Found {len(failed_claims)} failed claims to retry")
                    logger.info(f"  ✅ {len(successful_claims)} claims already verified successfully")
                else:
                    logger.info(f"  ✅ All {len(successful_claims)} claims already verified successfully")
                    return EvalResult(
                        metric_name=self.metric_name,
                        score=existing_data.get('score', 0),
                        details=details,
                        weight=self.weight
                    )
        
        phase1_time = 0
        if claims is None:
            logger.info("=" * 60)
            logger.info("Phase 1: Extracting fact-citation pairs using LLM...")
            logger.info(f"  Report length: {len(result_text)} chars")
            phase1_start = time.time()
            claims = self._extract_claims_with_llm(result_text)
            phase1_time = time.time() - phase1_start
            logger.info(f"  Phase 1 completed in {phase1_time:.1f}s")
            
            if not claims:
                logger.warning("  ⚠️ No claims extracted, returning error score -1")
                return EvalResult(metric_name=self.metric_name, score=-1,
                                details={'method': 'llm_agent_extraction', 'error': 'Failed to extract claims', 'claim_count': 0},
                                weight=self.weight)
            
            logger.info(f"  ✅ Extracted {len(claims)} fact-citation pairs")
        else:
            logger.info("  ⏭️ Skipping claim extraction (using existing claims)")
        
        logger.info("=" * 60)
        logger.info("Phase 2: Verifying claims...")
        phase2_start = time.time()
        
        if failed_claims:
            logger.info(f"  🔄 Only verifying {len(failed_claims)} failed claims...")
            new_verification_results = self._verify_claims_batch(failed_claims, source_folder)
            
            verification_results, accuracy_score = self._merge_verification_results(
                successful_claims, new_verification_results, claims
            )
            logger.info(f"  ✅ Merged {len(successful_claims)} successful + {len(new_verification_results)} new verifications")
        else:
            verification_results = self._verify_claims_batch(claims, source_folder)
        
        phase2_time = time.time() - phase2_start
        
        supported_count = sum(1 for r in verification_results if r.get('supported'))
        total_claims = len(claims)
        accuracy_score = (supported_count / total_claims * 100) if total_claims > 0 else 0
        
        long_context_claims = []
        source_claims = []
        
        for i, claim in enumerate(claims):
            source_type = claim.get('source_type', 'unknown')
            verification = verification_results[i] if i < len(verification_results) else {}
            
            if source_type in ('long_context', 'rag'):
                long_context_claims.append({
                    'claim': claim,
                    'verification': verification
                })
            else:
                source_claims.append({
                    'claim': claim,
                    'verification': verification
                })
        
        lc_supported = sum(1 for c in long_context_claims if c['verification'].get('supported'))
        lc_total = len(long_context_claims)
        lc_score = (lc_supported / lc_total * 100) if lc_total > 0 else -1
        
        src_supported = sum(1 for c in source_claims if c['verification'].get('supported'))
        src_total = len(source_claims)
        src_score = (src_supported / src_total * 100) if src_total > 0 else -1
        
        details = {
            'method': 'llm_agent_extraction',
            'extraction': {'total_claims': total_claims, 'claims': claims},
            'verification_result': {
                'verifications': verification_results,
                'supported_count': supported_count,
                'unsupported_count': total_claims - supported_count,
                'accuracy_score': accuracy_score
            },
            'score_breakdown': {
                'long_context': {
                    'score': lc_score,
                    'supported_count': lc_supported,
                    'total_count': lc_total,
                    'claims': [c['claim'].get('id') for c in long_context_claims]
                },
                'source': {
                    'score': src_score,
                    'supported_count': src_supported,
                    'total_count': src_total,
                    'claims': [c['claim'].get('id') for c in source_claims]
                }
            }
        }
        
        total_time = time.time() - eval_start_time
        lc_score_str = f"{lc_score:.1f}" if lc_score >= 0 else "N/A"
        src_score_str = f"{src_score:.1f}" if src_score >= 0 else "N/A"
        logger.info("=" * 60)
        logger.info(f"📊 EVALUATION COMPLETE")
        logger.info(f"   Final Score: {accuracy_score:.1f}/100")
        logger.info(f"   Supported: {supported_count}/{total_claims}")
        logger.info(f"   - Long Context: {lc_score_str}/100 ({lc_supported}/{lc_total})")
        logger.info(f"   - Source Files: {src_score_str}/100 ({src_supported}/{src_total})")
        logger.info(f"   Phase 1 (Extract): {phase1_time:.1f}s")
        logger.info(f"   Phase 2 (Verify):  {phase2_time:.1f}s")
        logger.info(f"   Total Time: {total_time:.1f}s")
        logger.info("=" * 60)
        return EvalResult(metric_name=self.metric_name, score=accuracy_score, details=details, weight=self.weight)
    
    def _extract_claims_with_llm(self, report_text: str) -> List[Dict]:
        """
        
        """
        if not self.api_client:
            logger.error("  ❌ API client not initialized")
            return []
        
        logger.info(f"  📊 Report length: {len(report_text)} chars (~{len(report_text)//4} tokens)")
        
        MAX_REPORT_LENGTH = 50000
        if len(report_text) > MAX_REPORT_LENGTH:
            logger.warning(f"  ⚠️ Report too long, truncating to {MAX_REPORT_LENGTH} chars")
            report_text = report_text[:MAX_REPORT_LENGTH]
        
        claims = self._extract_claims_single_call(report_text)
        
        if claims:
            for i, claim in enumerate(claims):
                claim['id'] = i + 1
        
        logger.info(f"  📊 Total claims extracted: {len(claims)}")
        return claims
    
    def _extract_claims_single_call(self, report_text: str) -> List[Dict]:
        prompt = f'''
你的任务是从报告中提取**所有**带引用的事实陈述，用于**验证事实准确性**。

我们需要检查报告中的每个事实陈述是否被其引用的源文档所支持。

{report_text}

引用标记是标注信息来源的标记，通常用方括号 `[...]` 包裹，但有时也会用小括号 `(...)` 代替。
**重要：任何方括号 `[...]` 内的内容都可能是引用标记，需要仔细识别。**

- `[long_context: "某文档", chunk 1]` 或 `[long_context: "某文档", chunk 1-3]`
- `[Doc: 报告.pdf]` 或 `[Doc: 报告.pdf, Page 5]` 或 `[Doc: 报告.pdf, Page 5-10]`
- `[文档: 报告.pdf]` 或 `[文档: 报告.pdf, 第5页]`
- `[RAG-1]` 或 `[RAG-1, RAG-2]`
- `[视频: xxx.mp4]` 或 `[Video: xxx.mp4]`
- `[图片: xxx.png]` 或 `[Image: xxx.png]`
- `[音频: xxx.mp3]` 或 `[Audio: xxx.mp3]`

有些模型会直接使用文件名作为引用，可能带有位置信息：
- `[llama.pdf]` → 文档引用
- `[llama.pdf, Section 7]` → 文档引用（带章节信息）
- `[report.pdf, Page 10-15]` → 文档引用（带页码范围）
- `[image.jpg]` → 图片引用
- `[video.mp4, 00:15:30]` → 视频引用（带时间戳）

有些模型会使用自定义的引用格式，例如：
- `[vision_understanding_advanced: image2.jpg, reasoning]` → 图片引用
- `[multimodal_analysis: video.mp4, scene 3]` → 视频引用
- `[document_reader: report.pdf, section 2]` → 文档引用
- `[source_1: data.csv]` → 数据文件引用

**识别规则：只要方括号内包含文件名（带扩展名如 .pdf, .jpg, .mp4 等），就应该识别为引用。根据文件扩展名确定 source_type：**
- `.pdf`, `.docx`, `.doc`, `.txt`, `.md`, `.html` → source_type: "doc"
- `.jpg`, `.jpeg`, `.png`, `.gif`, `.webp` → source_type: "image"
- `.mp4`, `.webm`, `.mov`, `.avi` → source_type: "video"
- `.mp3`, `.wav`, `.m4a`, `.aac`, `.ogg`, `.flac` → source_type: "audio"

有些引用会包含具体的页码或位置信息，请完整保留：
- `[Doc: llama.pdf, Page 34]` → 保留页码信息
- `[Doc: report.pdf, Page 10-15]` → 保留页码范围
- `[long_context: "论文", chunk 5-8]` → 保留 chunk 范围
- `[视频: demo.mp4, 00:15:30]` → 保留时间戳

有些报告会错误地使用小括号代替方括号，例如：
- `(long_context: "某文档", chunk 1)` → 等同于 `[long_context: "某文档", chunk 1]`
- `(Doc: 报告.pdf)` → 等同于 `[Doc: 报告.pdf]`
- `(RAG-1)` → 等同于 `[RAG-1]`

**请同时识别方括号和小括号格式的引用。**

1. **逐个扫描**：从报告开头到结尾，逐个找出每个引用标记（方括号或小括号内的内容）
2. **提取前置句子**：对于每个引用标记，提取它**紧邻前面**的完整句子作为事实陈述
3. **句子边界**：句子以句号、问号、感叹号、冒号结尾，或以换行符分隔
4. **多引用合并**：如果一个句子后面有多个连续的引用标记，将它们合并为一条记录
5. **去除引用标记**：claim_text 中不要包含引用标记本身
6. **保留完整引用**：citations 中保留引用标记的完整内容，包括页码、chunk 等位置信息

1. **参考文献/引用来源部分**：报告末尾列出的参考文献列表
2. **纯引用标记**：没有前置事实陈述的孤立引用标记
3. **标题和目录**：章节标题、目录项
4. **元信息**：如"本报告基于以下来源"等说明性文字
5. **Markdown 链接**：如 `[链接文字](url)` 这种格式不是引用

**输入报告片段**：
```
根据调查数据，2023年全球电动汽车销量达到1400万辆[Doc: 电动汽车报告.pdf, Page 12]。
其中中国市场占比超过60%[long_context: "市场分析", chunk 3][RAG-2]。
该公司的营收增长了25%(Doc: 财报.pdf)。
图片显示了产品的外观设计[vision_understanding_advanced: product.jpg, design_analysis]。
```

**正确提取结果**：
```json
{{
    "claims": [
        {{
            "id": 1,
            "claim_text": "根据调查数据，2023年全球电动汽车销量达到1400万辆",
            "citations": ["[Doc: 电动汽车报告.pdf, Page 12]"],
            "source_type": "doc"
        }},
        {{
            "id": 2,
            "claim_text": "其中中国市场占比超过60%",
            "citations": ["[long_context: \\"市场分析\\", chunk 3]", "[RAG-2]"],
            "source_type": "long_context"
        }},
        {{
            "id": 3,
            "claim_text": "该公司的营收增长了25%",
            "citations": ["[Doc: 财报.pdf]"],
            "source_type": "doc"
        }},
        {{
            "id": 4,
            "claim_text": "图片显示了产品的外观设计",
            "citations": ["[vision_understanding_advanced: product.jpg, design_analysis]"],
            "source_type": "image"
        }}
    ]
}}
```

{{
    "claims": [
        {{
            "id": 1,
            "claim_text": "事实陈述（不含引用标记）",
            "citations": ["[标准化引用格式]"],
            "source_type": "long_context 或 doc 或 rag 或 image 或 video 或 audio"
        }}
    ]
}}

**无论原始报告中使用什么格式，输出时必须统一转换为以下标准格式：**

| 原始格式示例 | 标准化输出格式 |
|-------------|---------------|
| `[Doc: xxx.pdf]`, `[doc: xxx.pdf]`, `[文档: xxx.pdf]`, `[xxx.pdf]` | `[Doc: xxx.pdf]` |
| `[Doc: xxx.pdf, Page 5]`, `[xxx.pdf, Section 7]` | `[Doc: xxx.pdf, Page 5]` |
| `[Image: xxx.jpg]`, `[image: xxx.jpg]`, `[图片: xxx.jpg]`, `[xxx.jpg]` | `[Image: xxx.jpg]` |
| `[Video: xxx.mp4]`, `[video: xxx.mp4]`, `[视频: xxx.mp4]`, `[xxx.mp4]` | `[Video: xxx.mp4]` |
| `[Audio: xxx.mp3]`, `[audio: xxx.mp3]`, `[音频: xxx.mp3]`, `[xxx.mp3]` | `[Audio: xxx.mp3]` |
| `[long_context: "xxx", chunk 1]` | `[long_context: "xxx", chunk 1]`（保持不变） |
| `[RAG-1]` | `[RAG-1]`（保持不变） |
| `[vision_understanding_advanced: xxx.jpg, ...]` | `[Image: xxx.jpg]` |
| `[multimodal_analysis: xxx.mp4, ...]` | `[Video: xxx.mp4]` |
| `(Doc: xxx.pdf)` 小括号格式 | `[Doc: xxx.pdf]` 转为方括号 |

**标准格式规则：**
1. **文档类型**：统一使用 `[Doc: 文件名]` 或 `[Doc: 文件名, Page X]`
2. **图片类型**：统一使用 `[Image: 文件名]`
3. **视频类型**：统一使用 `[Video: 文件名]` 或 `[Video: 文件名, 时间戳]`
4. **音频类型**：统一使用 `[Audio: 文件名]`
5. **long_context**：保持 `[long_context: "标题", chunk X]` 格式不变
6. **RAG**：保持 `[RAG-N]` 格式不变
7. **所有引用必须使用方括号 `[]`，不使用小括号 `()`**

1. **claim_text 中不要包含任何引号**：如果原文中有引号（单引号、双引号、中文引号），请用其他方式表达或直接省略
   - ❌ 错误：`"claim_text": "他说\"这很重要\""`
   - ✅ 正确：`"claim_text": "他说这很重要"`
   - ❌ 错误：`"claim_text": "所谓'创新'是指..."`
   - ✅ 正确：`"claim_text": "所谓创新是指..."`
2. **必须使用双引号**：JSON 的键和字符串值必须使用双引号 `"`，不能使用单引号 `'`
3. **确保 JSON 格式正确**：不要有多余的逗号、缺少逗号等语法错误

- 确保提取完整，不要遗漏任何带引用的事实陈述
- **必须将所有引用转换为上述标准格式**
- source_type 根据第一个引用标记的类型确定
- 保留引用中的页码、chunk、时间戳等位置信息
- **claim_text 中绝对不要包含任何引号字符**

请开始提取：'''

        for retry in range(API_MAX_RETRIES):
            try:
                logger.info(f"  📤 Sending request to LLM (attempt {retry + 1}/{API_MAX_RETRIES})...")
                
                max_tokens = get_max_tokens_for_model(MODEL_NAME)
                logger.info(f"  Using max_tokens={max_tokens} for model {MODEL_NAME}")
                
                response = self.api_client.chat.completions.create(
                    model=MODEL_NAME,
                    messages=[{"role": "system", "content": "你是一个专业的事实核查专家。"},
                              {"role": "user", "content": prompt}],
                    max_tokens=max_tokens, temperature=0)
                
                response_content = response.choices[0].message.content if response.choices else ""
                finish_reason = response.choices[0].finish_reason if response.choices else None
                
                if finish_reason == 'length':
                    logger.error(f"  ❌ Response truncated due to max_tokens limit!")
                    logger.error(f"  Completion tokens: {response.usage.completion_tokens if hasattr(response, 'usage') else 'N/A'}")
                    logger.error(f"  This report has too many claims. Consider splitting or reducing report length.")
                    if not response_content:
                        logger.error(f"  Content is empty even though tokens were used. This is unusual.")
                        return []
                
                if not response_content:
                    logger.warning(f"  ⚠️ Empty response from LLM (attempt {retry + 1}/{API_MAX_RETRIES})")
                    logger.warning(f"  Finish reason: {finish_reason}")
                    logger.warning(f"  Response object: {response}")
                    
                    if retry < API_MAX_RETRIES - 1:
                        logger.info(f"  ⏳ Retrying after {API_RETRY_DELAY}s...")
                        time.sleep(API_RETRY_DELAY)
                        continue
                    else:
                        logger.error(f"  ❌ Failed to get response after {API_MAX_RETRIES} attempts")
                        return []
                
                logger.info(f"  📥 Received response: {len(response_content)} chars")
                logger.debug(f"  Response preview: {response_content[:500]}...")
                
                if response_content:
                    result = extract_json_from_text(response_content)
                    
                    if result is not None:
                        claims = result.get('claims', [])
                        if claims:
                            logger.info(f"  ✅ Successfully extracted {len(claims)} claims")
                            return claims
                        else:
                            logger.warning(f"  ⚠️ LLM returned empty claims list")
                            logger.warning(f"  Parsed result: {result}")
                    else:
                        logger.warning(f"  ⚠️ extract_json_from_text failed")
                        logger.warning(f"  Raw response (first 1000 chars): {response_content[:1000]}")
                        
                        try:
                            cleaned = re.sub(r'```json\s*', '', response_content)
                            cleaned = re.sub(r'```\s*', '', cleaned)
                            
                            cleaned = cleaned.replace('"', '"').replace('"', '"')
                            cleaned = cleaned.replace(''', "'").replace(''', "'")
                            
                            json_start = cleaned.find("{")
                            if json_start >= 0:
                                json_str = cleaned[json_start:]
                                
                                try:
                                    import ast
                                    result = ast.literal_eval(json_str)
                                    claims = result.get('claims', [])
                                    if claims:
                                        logger.info(f"  ✅ Extracted {len(claims)} claims via ast.literal_eval")
                                        return claims
                                except (ValueError, SyntaxError):
                                    pass
                                
                                open_braces = json_str.count('{') - json_str.count('}')
                                open_brackets = json_str.count('[') - json_str.count(']')
                                
                                if open_braces > 0 or open_brackets > 0:
                                    logger.info(f"  🔧 Attempting to fix truncated JSON (missing {open_brackets} ] and {open_braces} }})")
                                    
                                    last_complete = -1
                                    for match in re.finditer(r'\}\s*,', json_str):
                                        last_complete = match.end()
                                    
                                    if last_complete > 0:
                                        truncated = json_str[:last_complete - 1]
                                        truncated += ']}'
                                        try:
                                            result = json.loads(truncated)
                                            claims = result.get('claims', [])
                                            if claims:
                                                logger.info(f"  ✅ Fixed truncated JSON, extracted {len(claims)} claims")
                                                return claims
                                        except json.JSONDecodeError:
                                            pass
                                    
                                    fixed_json = json_str
                                    if fixed_json.count('"') % 2 == 1:
                                        fixed_json += '"'
                                    fixed_json += ']' * open_brackets + '}' * open_braces
                                    try:
                                        result = json.loads(fixed_json)
                                        claims = result.get('claims', [])
                                        if claims:
                                            logger.info(f"  ✅ Fixed truncated JSON by adding brackets, extracted {len(claims)} claims")
                                            return claims
                                    except json.JSONDecodeError as e:
                                        logger.warning(f"  ⚠️ Simple fix failed: {e}")
                        except Exception as e:
                            logger.warning(f"  ⚠️ Fallback parsing failed: {e}")
                    
                    return []
                        
                return []
                
            except Exception as e:
                logger.error(f"  ❌ API error (attempt {retry + 1}/{API_MAX_RETRIES}): {e}")
                
                if '429' in str(e).lower() or 'rate' in str(e).lower():
                    if retry < API_MAX_RETRIES - 1:
                        logger.info(f"  ⏳ Rate limited, waiting {API_RETRY_DELAY}s before retry...")
                        time.sleep(API_RETRY_DELAY)
                        continue
                
                import traceback
                logger.error(f"  Stack trace: {traceback.format_exc()}")
                raise RuntimeError(f"API connection error: {e}")
        
        logger.error(f"  ❌ Failed after {API_MAX_RETRIES} attempts")
        raise RuntimeError(f"Failed to extract claims after {API_MAX_RETRIES} attempts")
    
    def _verify_claims_batch(self, claims: List[Dict], source_folder: Optional[Path]) -> List[Dict]:
        import time
        phase2_start = time.time()
        
        if not self.api_client:
            return [{'id': c.get('id'), 'supported': False, 'explanation': 'API not initialized'} for c in claims]
        
        verification_tasks = []
        for claim in claims:
            claim_id = claim.get('id')
            claim_text = claim.get('claim_text', '')
            citations = claim.get('citations', [])
            source_type = claim.get('source_type', 'unknown')
            
            if not citations:
                verification_tasks.append({'claim_id': claim_id, 'claim_text': claim_text, 'citation': None, 'source_key': 'unknown'})
            else:
                for citation in citations:
                    source_key = self._get_source_key([citation], source_type)
                    verification_tasks.append({'claim_id': claim_id, 'claim_text': claim_text, 'citation': citation, 'source_key': source_key})
        
        logger.info(f"  📋 Created {len(verification_tasks)} verification tasks for {len(claims)} claims")
        
        source_groups: Dict[str, List[Dict]] = {}
        for task in verification_tasks:
            source_key = task['source_key']
            if source_key not in source_groups:
                source_groups[source_key] = []
            source_groups[source_key].append(task)
        
        logger.info(f"  Grouped into {len(source_groups)} source documents")
        
        task_results: Dict[Tuple[int, str], Dict] = {}
        
        batch_jobs = []
        for source_key, group_tasks in source_groups.items():
            source_info = self._get_source_info(source_key, source_folder)
            for i in range(0, len(group_tasks), BATCH_SIZE):
                batch = group_tasks[i:i + BATCH_SIZE]
                batch_jobs.append((source_key, source_info, batch))
        
        logger.info(f"  Created {len(batch_jobs)} batch jobs (max {BATCH_SIZE} claims each)")
        
        def verify_single_batch(job):
            source_key, source_info, batch = job
            temp_claims = [{'id': idx + 1, 'claim_text': t['claim_text'], 'citation': t['citation']} for idx, t in enumerate(batch)]
            
            thread_client = OpenAI(api_key=API_KEY, base_url=API_BASE_URL)
            
            source_type = source_info.get('type', '')
            file_path = source_info.get('file_path')
            file_ext = file_path.suffix.lower() if file_path else ''
            
            model_supports_image = 'gpt-5' in MODEL_NAME.lower() or 'gpt5' in MODEL_NAME.lower()
            
            if source_type == 'audio':
                results = self._verify_audio_with_gemini_thread(temp_claims, source_info, thread_client)
            elif source_type == 'video':
                results = self._verify_video_with_gemini_thread(temp_claims, source_info, thread_client)
            elif source_type == 'file' and file_ext in IMAGE_EXTENSIONS:
                if model_supports_image:
                    results = self._verify_batch_with_api_thread(temp_claims, source_info, thread_client)
                else:
                    results = self._verify_image_with_gemini_thread(temp_claims, source_info, thread_client)
            else:
                results = self._verify_batch_with_api_thread(temp_claims, source_info, thread_client)
            
            batch_results = []
            for idx, task in enumerate(batch):
                result = results[idx] if idx < len(results) else {'supported': False, 'explanation': 'No result'}
                task_key = (task['claim_id'], task['citation'] or '')
                batch_results.append((task_key, {
                    'supported': result.get('supported', False),
                    'explanation': result.get('explanation', ''),
                    'source_found': result.get('source_found', source_info.get('found', False))
                }))
            return source_key, batch_results
        
        max_workers = min(len(batch_jobs), MAX_CONCURRENT_REQUESTS)
        logger.info(f"  Starting parallel verification with {max_workers} workers (max {MAX_CONCURRENT_REQUESTS})")
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(verify_single_batch, job): job for job in batch_jobs}
            batch_num = 0
            for future in as_completed(futures):
                batch_num += 1
                try:
                    source_key, batch_results = future.result()
                    logger.info(f"  Batch {batch_num}/{len(batch_jobs)}: Completed '{source_key[:30]}...'")
                    for task_key, result in batch_results:
                        task_results[task_key] = result
                        status = "✅" if result.get('supported') else "❌"
                        logger.info(f"    {status} [Claim {task_key[0]}]: {result.get('explanation', '')[:50]}")
                except Exception as e:
                    logger.error(f"  Batch {batch_num} failed: {e}")
        
        all_results = []
        for claim in claims:
            claim_id = claim.get('id')
            citations = claim.get('citations', [])
            
            if not citations:
                all_results.append({'id': claim_id, 'claim_text': claim.get('claim_text', '')[:100], 'supported': False, 'source_found': False, 'explanation': 'No citations', 'citation_results': []})
                continue
            
            citation_results = []
            all_supported = True
            any_source_found = False
            explanations = []
            
            for citation in citations:
                task_key = (claim_id, citation)
                result = task_results.get(task_key, {'supported': False, 'explanation': 'Not verified'})
                citation_results.append({'citation': citation, 'supported': result.get('supported', False), 'explanation': result.get('explanation', '')})
                if not result.get('supported', False):
                    all_supported = False
                if result.get('source_found', False):
                    any_source_found = True
                explanations.append(f"{citation}: {result.get('explanation', '')[:30]}")
            
            all_results.append({'id': claim_id, 'claim_text': claim.get('claim_text', '')[:100], 'supported': all_supported, 'source_found': any_source_found, 'explanation': '; '.join(explanations), 'citation_results': citation_results})
        
        return all_results
    
    def _get_source_key(self, citations: List[str], source_type: str) -> str:
        if not citations:
            return 'unknown'
        first_citation = citations[0]
        
        if 'long_context:' in first_citation:
            match = re.search(r'long_context:\s*"([^"]+)"', first_citation)
            if match:
                return f"long_context:{match.group(1)}"
        elif 'Doc:' in first_citation or '文档:' in first_citation or 'doc:' in first_citation.lower():
            match = re.search(r'(?:[Dd]oc|文档):\s*([^,\]]+)', first_citation)
            if match:
                return f"doc:{match.group(1).strip()}"
        elif 'RAG-' in first_citation:
            match = re.search(r'RAG-(\d+)', first_citation)
            if match:
                return f"rag:{match.group(1)}"
        elif '图片:' in first_citation or 'Image:' in first_citation.lower():
            match = re.search(r'(?:图片|[Ii]mage):\s*([^,\]]+)', first_citation)
            if match:
                return f"image:{match.group(1).strip()}"
        elif '视频:' in first_citation or 'video:' in first_citation.lower():
            match = re.search(r'(?:视频|[Vv]ideo):\s*([^,\]]+)', first_citation)
            if match:
                return f"video:{match.group(1).strip()}"
        elif '音频:' in first_citation or 'audio:' in first_citation.lower():
            match = re.search(r'(?:音频|[Aa]udio):\s*([^,\]]+)', first_citation)
            if match:
                return f"audio:{match.group(1).strip()}"
        doc_match = re.search(r'([^\[\],\s]+\.(?:pdf|md|docx?|txt))', first_citation, re.I)
        if doc_match:
            return f"doc:{doc_match.group(1).strip()}"
        img_match = re.search(r'([^\[\],\s]+\.(?:jpg|jpeg|png|gif|webp))', first_citation, re.I)
        if img_match:
            return f"image:{img_match.group(1).strip()}"
        video_match = re.search(r'([^\[\],\s]+\.(?:mp4|webm|mov|avi))', first_citation, re.I)
        if video_match:
            return f"video:{video_match.group(1).strip()}"
        audio_match = re.search(r'([^\[\],\s]+\.(?:mp3|wav|m4a|aac|ogg|flac|wma))', first_citation, re.I)
        if audio_match:
            return f"audio:{audio_match.group(1).strip()}"
        return f"unknown:{first_citation[:30]}"
    
    def _get_source_info(self, source_key: str, source_folder: Optional[Path]) -> Dict:
        parts = source_key.split(':', 1)
        source_type = parts[0] if len(parts) > 1 else 'unknown'
        source_name = parts[1] if len(parts) > 1 else source_key
        
        if source_type == 'long_context':
            content, name = self.doc_loader.get_content_by_title(source_name)
            if content:
                return {'found': True, 'type': 'text', 'content': content, 'name': name}
        elif source_type == 'doc':
            if source_folder:
                for folder in [source_folder / "source", source_folder]:
                    if not folder.exists():
                        continue
                    for file_path in folder.iterdir():
                        if file_path.is_file() and fuzzy_title_match(source_name, file_path.name):
                            return {'found': True, 'type': 'file', 'file_path': file_path, 'name': file_path.name}
            for doc_name, content in self.doc_loader.source_documents.items():
                if fuzzy_title_match(source_name, doc_name):
                    return {'found': True, 'type': 'text', 'content': content, 'name': doc_name}
            content, name = self.doc_loader.get_content_by_title(source_name)
            if content:
                return {'found': True, 'type': 'text', 'content': content, 'name': name}
        elif source_type == 'rag':
            try:
                idx = int(source_name) - 1
                content, name = self.doc_loader.get_content_by_index(idx)
                if content:
                    return {'found': True, 'type': 'text', 'content': content, 'name': name}
            except ValueError:
                pass
        elif source_type == 'image':
            if source_folder:
                for folder in [source_folder / "source", source_folder]:
                    if not folder.exists():
                        continue
                    for file_path in folder.iterdir():
                        if file_path.is_file() and file_path.suffix.lower() in IMAGE_EXTENSIONS:
                            if fuzzy_title_match(source_name, file_path.name):
                                return {'found': True, 'type': 'file', 'file_path': file_path, 'name': file_path.name}
        elif source_type == 'video':
            if source_folder:
                for folder in [source_folder / "source", source_folder]:
                    if not folder.exists():
                        continue
                    for file_path in folder.iterdir():
                        if file_path.is_file() and file_path.suffix.lower() in VIDEO_EXTENSIONS:
                            if fuzzy_title_match(source_name, file_path.name):
                                logger.info(f"    ✅ Found video: {file_path.name}")
                                return {'found': True, 'type': 'video', 'file_path': file_path, 'name': file_path.name}
        elif source_type == 'audio':
            if source_folder:
                for folder in [source_folder / "source", source_folder]:
                    if not folder.exists():
                        continue
                    for file_path in folder.iterdir():
                        if file_path.is_file() and file_path.suffix.lower() in AUDIO_EXTENSIONS:
                            if fuzzy_title_match(source_name, file_path.name):
                                logger.info(f"    ✅ Found audio: {file_path.name}")
                                return {'found': True, 'type': 'audio', 'file_path': file_path, 'name': file_path.name}
        return {'found': False, 'type': 'not_found', 'name': source_name}
    
    def _verify_audio_with_gemini_thread(self, claims: List[Dict], source_info: Dict, client: OpenAI) -> List[Dict]:
        if not source_info.get('found'):
            return [{'id': c.get('id'), 'supported': False, 'source_found': False, 'explanation': 'Audio not found'} for c in claims]
        
        file_path = source_info['file_path']
        audio_content, description = prepare_audio_for_gemini(file_path)
        
        if not audio_content:
            return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Cannot prepare audio'} for c in claims]
        
        claims_text = "\n".join([f"[{c.get('id')}] {c.get('claim_text')}" for c in claims])
        prompt = f"""请验证以下陈述是否被音频内容支持。

**待验证的陈述**:
{claims_text}

**验证标准（宽松）**:
- 只要陈述中有任何内容在音频中被提及，就判定为 supported: true
- 允许概括和推断
- 只有直接矛盾或完全无中生有才判定为 false

请以 JSON 格式回复:
{{"results": [{{"id": 1, "supported": true, "explanation": "简要解释"}}]}}"""
        
        content_parts = audio_content + [{"type": "text", "text": prompt}]
        
        for retry in range(API_MAX_RETRIES):
            try:
                response = client.chat.completions.create(
                    model=GEMINI_MULTIMODAL_MODEL,
                    messages=[{"role": "user", "content": content_parts}],
                    max_tokens=4096
                )
                response_content = response.choices[0].message.content if response.choices else ""
                if response_content:
                    try:
                        cleaned = re.sub(r'```json\s*', '', response_content)
                        cleaned = re.sub(r'```\s*', '', cleaned)
                        json_start = cleaned.find("{")
                        json_end = cleaned.rfind("}") + 1
                        if json_start >= 0 and json_end > json_start:
                            result_data = json.loads(cleaned[json_start:json_end])
                            result_map = {r.get('id'): r for r in result_data.get('results', [])}
                            return [{'id': c.get('id'), 'supported': result_map.get(c.get('id'), {}).get('supported', False), 'source_found': True, 'explanation': result_map.get(c.get('id'), {}).get('explanation', '')} for c in claims]
                    except json.JSONDecodeError as e:
                        logger.error(f"  ❌ Audio JSON parse error: {e}")
                        logger.error(f"  Raw response (first 500 chars): {response_content[:500]}")
                else:
                    logger.warning(f"  ⚠️ Empty response from Gemini for audio verification")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'Parse error: response={response_content[:100] if response_content else "empty"}'} for c in claims]
            except Exception as e:
                if '429' in str(e).lower() or 'rate' in str(e).lower():
                    if retry < API_MAX_RETRIES - 1:
                        time.sleep(API_RETRY_DELAY)
                        continue
                logger.error(f"  ❌ Gemini API error: {e}")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'API error: {e}'} for c in claims]
        return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Max retries exceeded'} for c in claims]
    
    def _verify_audio_with_gemini(self, claims: List[Dict], source_info: Dict) -> List[Dict]:
        if not self.gemini_client:
            return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Gemini client not available'} for c in claims]
        return self._verify_audio_with_gemini_thread(claims, source_info, self.gemini_client)
    
    def _verify_video_with_gemini_thread(self, claims: List[Dict], source_info: Dict, client: OpenAI) -> List[Dict]:
        """Verify claims against video content using Gemini multimodal model."""
        if not source_info.get('found'):
            return [{'id': c.get('id'), 'supported': False, 'source_found': False, 'explanation': 'Video not found'} for c in claims]
        
        file_path = source_info['file_path']
        
        video_content, description = prepare_video_for_gemini(file_path)
        
        if not video_content:
            return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Cannot prepare video'} for c in claims]
        
        claims_text = "\n".join([f"[{c.get('id')}] {c.get('claim_text')}" for c in claims])
        
        prompt_full = f"""请验证以下陈述是否被视频内容支持。

**待验证的陈述**:
{claims_text}

**验证标准（宽松）**:
- 只要陈述中有任何内容在视频中被提及或展示，就判定为 supported: true
- 允许概括和推断
- 只有直接矛盾或完全无中生有才判定为 false

请以 JSON 格式回复:
{{"results": [{{"id": 1, "supported": true, "explanation": "简要解释"}}]}}"""
        
        prompt_simple = f"""请验证以下陈述是否被视频内容支持。只需回复 true 或 false。

**待验证的陈述**:
{claims_text}

**验证标准（宽松）**:
- 只要陈述中有任何内容在视频中被提及或展示，就判定为 true
- 允许概括和推断
- 只有直接矛盾或完全无中生有才判定为 false

请以 JSON 格式回复（不需要解释）:
{{"results": [{{"id": 1, "supported": true}}]}}"""
        
        use_simple_prompt = False
        
        for retry in range(API_MAX_RETRIES):
            try:
                current_prompt = prompt_simple if use_simple_prompt else prompt_full
                content_parts = video_content + [{"type": "text", "text": current_prompt}]
                
                response = client.chat.completions.create(
                    model=GEMINI_MULTIMODAL_MODEL,
                    messages=[{"role": "user", "content": content_parts}],
                    max_tokens=4096
                )
                response_content = response.choices[0].message.content if response.choices else ""
                
                if response_content:
                    try:
                        cleaned = re.sub(r'```json\s*', '', response_content)
                        cleaned = re.sub(r'```\s*', '', cleaned)
                        json_start = cleaned.find("{")
                        json_end = cleaned.rfind("}") + 1
                        if json_start >= 0 and json_end > json_start:
                            json_str = cleaned[json_start:json_end]
                            try:
                                result_data = json.loads(json_str)
                            except json.JSONDecodeError:
                                logger.warning(f"  ⚠️ Video JSON truncated, attempting to fix...")
                                fixed_json = json_str
                                open_braces = fixed_json.count('{') - fixed_json.count('}')
                                open_brackets = fixed_json.count('[') - fixed_json.count(']')
                                if fixed_json.count('"') % 2 == 1:
                                    fixed_json += '"'
                                fixed_json += ']' * open_brackets + '}' * open_braces
                                try:
                                    result_data = json.loads(fixed_json)
                                    logger.info(f"  ✅ Fixed truncated video JSON by adding {open_brackets} ] and {open_braces} }}")
                                except json.JSONDecodeError as e2:
                                    if not use_simple_prompt and retry < API_MAX_RETRIES - 1:
                                        logger.warning(f"  ⚠️ JSON fix failed, switching to simple prompt for retry...")
                                        use_simple_prompt = True
                                        time.sleep(API_RETRY_DELAY)
                                        continue
                                    else:
                                        raise e2
                            
                            result_map = {r.get('id'): r for r in result_data.get('results', [])}
                            return [{'id': c.get('id'), 'supported': result_map.get(c.get('id'), {}).get('supported', False), 'source_found': True, 'explanation': result_map.get(c.get('id'), {}).get('explanation', 'verified via video')} for c in claims]
                    except json.JSONDecodeError as e:
                        logger.error(f"  ❌ Video JSON parse error: {e}")
                        logger.error(f"  Raw response (first 500 chars): {response_content[:500]}")
                        if not use_simple_prompt and retry < API_MAX_RETRIES - 1:
                            logger.warning(f"  ⚠️ Switching to simple prompt for retry...")
                            use_simple_prompt = True
                            time.sleep(API_RETRY_DELAY)
                            continue
                else:
                    logger.warning(f"  ⚠️ Empty response from Gemini for video verification")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'Parse error: response={response_content[:100] if response_content else "empty"}'} for c in claims]
            except Exception as e:
                if '429' in str(e).lower() or 'rate' in str(e).lower():
                    if retry < API_MAX_RETRIES - 1:
                        time.sleep(API_RETRY_DELAY)
                        continue
                logger.error(f"  ❌ Gemini API error (video): {e}")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'API error: {e}'} for c in claims]
        return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Max retries exceeded'} for c in claims]
    
    def _verify_video_with_gemini(self, claims: List[Dict], source_info: Dict) -> List[Dict]:
        if not self.gemini_client:
            return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Gemini client not available'} for c in claims]
        return self._verify_video_with_gemini_thread(claims, source_info, self.gemini_client)
    
    def _verify_image_with_gemini_thread(self, claims: List[Dict], source_info: Dict, client: OpenAI) -> List[Dict]:
        """Verify claims against image content using Gemini multimodal model."""
        if not source_info.get('found'):
            return [{'id': c.get('id'), 'supported': False, 'source_found': False, 'explanation': 'Image not found'} for c in claims]
        
        file_path = source_info['file_path']
        
        try:
            ext = file_path.suffix.lower()
            mime_types = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", 
                          ".gif": "image/gif", ".webp": "image/webp"}
            mime_type = mime_types.get(ext, "image/jpeg")
            image_base64 = encode_file_to_base64(str(file_path))
            image_content = [{"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_base64}"}}]
        except Exception as e:
            logger.error(f"  ❌ Cannot prepare image: {e}")
            return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'Cannot prepare image: {e}'} for c in claims]
        
        claims_text = "\n".join([f"[{c.get('id')}] {c.get('claim_text')}" for c in claims])
        prompt = f"""请验证以下陈述是否被图片内容支持。

**待验证的陈述**:
{claims_text}

**验证标准（宽松）**:
- 只要陈述中有任何内容在图片中可见或可推断，就判定为 supported: true
- 允许概括和推断
- 允许主观描述（如"设计精美"、"色彩丰富"等）
- 只有直接矛盾或完全无中生有才判定为 false

请以 JSON 格式回复:
{{"results": [{{"id": 1, "supported": true, "explanation": "简要解释"}}]}}"""
        
        content_parts = image_content + [{"type": "text", "text": prompt}]
        
        for retry in range(API_MAX_RETRIES):
            try:
                response = client.chat.completions.create(
                    model=GEMINI_MULTIMODAL_MODEL,
                    messages=[{"role": "user", "content": content_parts}],
                    max_tokens=4096
                )
                response_content = response.choices[0].message.content if response.choices else ""
                
                if response_content:
                    try:
                        cleaned = re.sub(r'```json\s*', '', response_content)
                        cleaned = re.sub(r'```\s*', '', cleaned)
                        json_start = cleaned.find("{")
                        json_end = cleaned.rfind("}") + 1
                        if json_start >= 0 and json_end > json_start:
                            result_data = json.loads(cleaned[json_start:json_end])
                            result_map = {r.get('id'): r for r in result_data.get('results', [])}
                            return [{'id': c.get('id'), 'supported': result_map.get(c.get('id'), {}).get('supported', False), 'source_found': True, 'explanation': result_map.get(c.get('id'), {}).get('explanation', '')} for c in claims]
                    except json.JSONDecodeError as e:
                        logger.error(f"  ❌ Image JSON parse error: {e}")
                        logger.error(f"  Raw response (first 500 chars): {response_content[:500]}")
                else:
                    logger.warning(f"  ⚠️ Empty response from Gemini for image verification")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'Parse error: response={response_content[:100] if response_content else "empty"}'} for c in claims]
            except Exception as e:
                if '429' in str(e).lower() or 'rate' in str(e).lower():
                    if retry < API_MAX_RETRIES - 1:
                        time.sleep(API_RETRY_DELAY)
                        continue
                logger.error(f"  ❌ Gemini API error (image): {e}")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'API error: {e}'} for c in claims]
        return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Max retries exceeded'} for c in claims]
    
    def _verify_pdf_with_gpt5_thread(self, claims: List[Dict], source_info: Dict, client: OpenAI) -> List[Dict]:
        """Verify claims against PDF content using GPT-5 responses API."""
        if not source_info.get('found'):
            return [{'id': c.get('id'), 'supported': False, 'source_found': False, 'explanation': 'PDF not found'} for c in claims]
        
        file_path = source_info['file_path']
        
        try:
            pdf_base64 = encode_file_to_base64(str(file_path))
            filename = file_path.name
            content_parts = [
                {"type": "input_file", "filename": filename, "file_data": f"data:application/pdf;base64,{pdf_base64}"}
            ]
            logger.info(f"    📄 Prepared PDF for GPT-5: {filename}")
        except Exception as e:
            logger.error(f"  ❌ Cannot prepare PDF: {e}")
            return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'Cannot prepare PDF: {e}'} for c in claims]
        
        claims_lines = []
        for c in claims:
            claim_id = c.get('id')
            claim_text = c.get('claim_text', '')
            citation = c.get('citation', '')
            
            location_hint = ''
            if citation:
                citation_details = parse_citation_details(citation)
                if citation_details.get('page'):
                    if citation_details.get('page_range') and citation_details['page_range'][0] != citation_details['page_range'][1]:
                        location_hint = f" (请重点关注第 {citation_details['page_range'][0]}-{citation_details['page_range'][1]} 页)"
                    else:
                        location_hint = f" (请重点关注第 {citation_details['page']} 页)"
            
            claims_lines.append(f"[{claim_id}] {claim_text}{location_hint}")
        
        claims_text = "\n".join(claims_lines)
        prompt = f"""
你是一个事实核查专家。请验证以下陈述是否与 PDF 文档内容**相关且不矛盾**。

**重要原则：保持评测的稳定，不要产生波动，评价标准较为宽松，只有明显错误才判 false。**

{claims_text}

只要满足以下**任意一条**，就判定为 true：
1. 陈述中的**任何部分**在 PDF 文档中有相关内容
2. 陈述是对 PDF 内容的**概括、总结、推断或延伸**
3. 陈述使用了不同的措辞，但**语义相近**
4. 陈述包含 PDF 中**隐含的信息**
5. 陈述是对 PDF 内容的**合理解读**，即使不是唯一解读
6. PDF 中有**部分支持**该陈述的内容

**只有**满足以下条件之一，才判定为 false：
1. 陈述与 PDF 内容**直接矛盾**（如数字存在明显误差、事实完全相反）
2. PDF 中**完全没有**陈述的任何相关内容
3. 陈述**不可能**从 PDF 合理推断出来

- 允许**大幅度的概括和推断**
- 允许**措辞差异**和**表述方式不同**
- 允许**部分正确**的陈述（只要不是完全错误）
- 对于**模糊或不确定**的情况，一律判定为 true

{{"results": [{{"id": 1, "supported": true, "explanation": "简要解释"}}]}}

请开始验证："""
        
        content_parts.append({"type": "input_text", "text": prompt})
        
        for retry in range(API_MAX_RETRIES):
            try:
                response = client.responses.create(
                    model=GPT5_MODEL,
                    input=[{"role": "user", "content": content_parts}]
                )
                
                response_content = ""
                if hasattr(response, 'output') and response.output:
                    for output in response.output:
                        if hasattr(output, 'content'):
                            for item in output.content:
                                if hasattr(item, 'text'):
                                    response_content += item.text
                
                if response_content:
                    try:
                        cleaned = re.sub(r'```json\s*', '', response_content)
                        cleaned = re.sub(r'```\s*', '', cleaned)
                        json_start = cleaned.find("{")
                        json_end = cleaned.rfind("}") + 1
                        if json_start >= 0 and json_end > json_start:
                            json_str = cleaned[json_start:json_end]
                            try:
                                result_data = json.loads(json_str)
                            except json.JSONDecodeError:
                                try:
                                    import ast
                                    result_data = ast.literal_eval(json_str)
                                except:
                                    fixed_json = json_str
                                    open_braces = fixed_json.count('{') - fixed_json.count('}')
                                    open_brackets = fixed_json.count('[') - fixed_json.count(']')
                                    fixed_json += ']' * open_brackets + '}' * open_braces
                                    result_data = json.loads(fixed_json)
                                    logger.info(f"  ✅ Fixed truncated JSON by adding {open_brackets} ] and {open_braces} }}")
                            
                            result_map = {r.get('id'): r for r in result_data.get('results', [])}
                            return [{'id': c.get('id'), 'supported': result_map.get(c.get('id'), {}).get('supported', False), 'source_found': True, 'explanation': result_map.get(c.get('id'), {}).get('explanation', '')} for c in claims]
                    except json.JSONDecodeError as e:
                        logger.error(f"  ❌ PDF JSON parse error: {e}")
                        logger.error(f"  Raw response (first 500 chars): {response_content[:500]}")
                else:
                    logger.warning(f"  ⚠️ Empty response from GPT-5 for PDF verification")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'Parse error: response={response_content[:100] if response_content else "empty"}'} for c in claims]
            except Exception as e:
                if '429' in str(e).lower() or 'rate' in str(e).lower():
                    if retry < API_MAX_RETRIES - 1:
                        time.sleep(API_RETRY_DELAY)
                        continue
                logger.error(f"  ❌ GPT-5 API error (PDF): {e}")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'API error: {e}'} for c in claims]
        return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Max retries exceeded'} for c in claims]
    
    def _verify_batch_with_api_thread(self, claims: List[Dict], source_info: Dict, client: OpenAI) -> List[Dict]:
        if not source_info.get('found'):
            return [{'id': c.get('id'), 'supported': False, 'source_found': False, 'explanation': 'Source not found'} for c in claims]
        
        content_parts = []
        api_type = 'chat'
        use_gemini = False
        
        model_supports_image = 'gpt-5' in MODEL_NAME.lower() or 'gpt5' in MODEL_NAME.lower()
        
        if source_info['type'] == 'file':
            file_path = source_info['file_path']
            file_ext = file_path.suffix.lower() if file_path else ''
            
            if file_ext in IMAGE_EXTENSIONS and not model_supports_image:
                use_gemini = True
            elif file_ext in VIDEO_EXTENSIONS:
                use_gemini = True
            
            if file_ext in PPTX_EXTENSIONS:
                page_numbers = get_page_numbers_from_claims(claims)
                logger.info(f"    📊 PPTX file detected, extracting pages: {page_numbers if page_numbers else 'all'}")
                text_content = extract_pptx_pages(file_path, page_numbers if page_numbers else None)
                if text_content:
                    if len(text_content) > MAX_CONTENT_CHARS:
                        logger.info(f"    🔍 PPTX content too long ({len(text_content)} chars), using BM25 retrieval...")
                        text_content = bm25_retrieve_for_claims(text_content, claims)
                    content_parts.append({"type": "text", "text": f"PPT 文档内容 ({file_path.name}):\n\n{text_content}"})
                else:
                    return [{'id': c.get('id'), 'supported': False, 'explanation': 'Cannot extract PPTX content'} for c in claims]
            elif file_ext in {'.docx', '.doc'}:
                logger.info(f"    📄 DOCX file detected, extracting content...")
                text_content = extract_docx_pages(file_path)
                if text_content:
                    if len(text_content) > MAX_CONTENT_CHARS:
                        logger.info(f"    🔍 DOCX content too long ({len(text_content)} chars), using BM25 retrieval...")
                        text_content = bm25_retrieve_for_claims(text_content, claims)
                    content_parts.append({"type": "text", "text": f"Word 文档内容 ({file_path.name}):\n\n{text_content}"})
                else:
                    return [{'id': c.get('id'), 'supported': False, 'explanation': 'Cannot extract DOCX content'} for c in claims]
            elif file_ext in {'.csv', '.tsv'}:
                logger.info(f"    📊 CSV/TSV file detected, extracting content...")
                text_content = extract_csv_content(file_path)
                if text_content:
                    if len(text_content) > MAX_CONTENT_CHARS:
                        logger.info(f"    🔍 CSV content too long ({len(text_content)} chars), using BM25 retrieval...")
                        text_content = bm25_retrieve_for_claims(text_content, claims)
                    content_parts.append({"type": "text", "text": f"CSV 数据文件内容 ({file_path.name}):\n\n{text_content}"})
                else:
                    return [{'id': c.get('id'), 'supported': False, 'explanation': 'Cannot extract CSV content'} for c in claims]
            elif file_ext in {'.xlsx', '.xls'}:
                logger.info(f"    📊 Excel file detected, extracting content...")
                text_content = extract_excel_content(file_path)
                if text_content:
                    if len(text_content) > MAX_CONTENT_CHARS:
                        logger.info(f"    🔍 Excel content too long ({len(text_content)} chars), using BM25 retrieval...")
                        text_content = bm25_retrieve_for_claims(text_content, claims)
                    content_parts.append({"type": "text", "text": f"Excel 数据文件内容 ({file_path.name}):\n\n{text_content}"})
                else:
                    return [{'id': c.get('id'), 'supported': False, 'explanation': 'Cannot extract Excel content'} for c in claims]
            elif file_ext in {'.json', '.jsonl'}:
                logger.info(f"    📊 JSON file detected, reading content...")
                try:
                    text_content = file_path.read_text(encoding='utf-8', errors='ignore')
                    if len(text_content) > MAX_CONTENT_CHARS:
                        logger.warning(f"    ⚠️ JSON content too long ({len(text_content)} chars), truncating to {MAX_CONTENT_CHARS}")
                        text_content = text_content[:MAX_CONTENT_CHARS] + "\n\n[... 内容已截断 ...]"
                    content_parts.append({"type": "text", "text": f"JSON 数据文件内容 ({file_path.name}):\n\n{text_content}"})
                except Exception as e:
                    return [{'id': c.get('id'), 'supported': False, 'explanation': f'Cannot read JSON: {e}'} for c in claims]
            else:
                file_content, description, api_type = prepare_file_for_api(file_path)
                if file_content:
                    content_parts.extend(file_content)
                else:
                    try:
                        text_content = file_path.read_text(encoding='utf-8', errors='ignore')
                        if len(text_content) > MAX_CONTENT_CHARS:
                            logger.warning(f"    ⚠️ File content too long ({len(text_content)} chars), truncating to {MAX_CONTENT_CHARS}")
                            text_content = text_content[:MAX_CONTENT_CHARS] + "\n\n[... 内容已截断 ...]"
                        content_parts.append({"type": "text", "text": f"文档内容 ({file_path.name}):\n\n{text_content}"})
                    except Exception:
                        return [{'id': c.get('id'), 'supported': False, 'explanation': 'Cannot read file'} for c in claims]
        else:
            text_content = source_info['content']
            if len(text_content) > MAX_CONTENT_CHARS:
                logger.info(f"    🔍 Text content too long ({len(text_content)} chars), using BM25 retrieval...")
                text_content = bm25_retrieve_for_claims(text_content, claims)
            content_parts.append({"type": "text", "text": f"文档内容 ({source_info['name']}):\n\n{text_content}"})
        
        claims_lines = []
        for c in claims:
            claim_id = c.get('id')
            claim_text = c.get('claim_text', '')
            citation = c.get('citation', '')
            
            location_hint = ''
            if citation:
                citation_details = parse_citation_details(citation)
                if citation_details.get('page'):
                    if citation_details.get('page_range') and citation_details['page_range'][0] != citation_details['page_range'][1]:
                        location_hint = f" (请重点关注第 {citation_details['page_range'][0]}-{citation_details['page_range'][1]} 页)"
                    else:
                        location_hint = f" (请重点关注第 {citation_details['page']} 页)"
                elif citation_details.get('chunk'):
                    if citation_details.get('chunk_range') and citation_details['chunk_range'][0] != citation_details['chunk_range'][1]:
                        location_hint = f" (请重点关注 chunk {citation_details['chunk_range'][0]}-{citation_details['chunk_range'][1]})"
                    else:
                        location_hint = f" (请重点关注 chunk {citation_details['chunk']})"
            
            claims_lines.append(f"[{claim_id}] {claim_text}{location_hint}")
        
        claims_text = "\n".join(claims_lines)
        prompt = f"""
你是一个事实核查专家。请验证以下陈述是否与源文档内容**相关且不矛盾**。

**重要原则：保持评测的稳定，不要产生波动，评价标准较为宽松，只有明显错误才判 false。**

{claims_text}

只要满足以下**任意一条**，就判定为 true：
1. 陈述中的**任何部分**在源文档中有相关内容
2. 陈述是对源文档内容的**概括、总结、推断或延伸**
3. 陈述使用了不同的措辞，但**语义相近**
4. 陈述包含源文档中**隐含的信息**
5. 对于图片/视频：陈述描述的内容**可能**在视觉上可见或可推断
6. 陈述是对源文档内容的**合理解读**，即使不是唯一解读
7. 源文档中有**部分支持**该陈述的内容

**只有**满足以下条件之一，才判定为 false：
1. 陈述与源文档**直接矛盾**（如数字存在明显误差、事实完全相反）
2. 源文档中**完全没有**陈述的任何相关内容
3. 陈述**不可能**从源文档合理推断出来

- 允许**大幅度的概括和推断**
- 允许**措辞差异**和**表述方式不同**
- 允许**部分正确**的陈述（只要不是完全错误）
- 对于**模糊或不确定**的情况，一律判定为 true
- 对于多模态内容，允许**主观描述**和**合理想象**

- 源文档："公司2023年营收增长了25%"
  - "公司营收大幅增长" → true（概括）
  - "公司业绩表现良好" → true（推断）
  - "公司2023年营收增长了约25%" → true（近似）
  - "公司2023年营收增长了约55%" → false（不符合事实）
  - "公司2023年营收下降了25%" → false（直接矛盾）

- 源文档：一张产品图片
  - "产品外观设计精美" → true（主观描述）
  - "产品采用了现代设计风格" → true（合理推断）

{{"results": [{{"id": 1, "supported": true, "explanation": "简要解释"}}]}}

请开始验证："""
        
        if api_type == 'responses':
            content_parts.append({"type": "input_text", "text": prompt})
        else:
            content_parts.append({"type": "text", "text": prompt})
        
        for retry in range(API_MAX_RETRIES):
            try:
                if api_type == 'responses':
                    response = client.responses.create(model=MODEL_NAME, input=[{"role": "user", "content": content_parts}])
                    response_content = ""
                    if hasattr(response, 'output') and response.output:
                        for output in response.output:
                            if hasattr(output, 'content'):
                                for item in output.content:
                                    if hasattr(item, 'text'):
                                        response_content += item.text
                else:
                    model_to_use = GEMINI_MULTIMODAL_MODEL if use_gemini else MODEL_NAME
                    response = client.chat.completions.create(model=model_to_use, messages=[{"role": "user", "content": content_parts}], max_tokens=4096)
                    response_content = response.choices[0].message.content if response.choices else ""
                
                if response_content:
                    try:
                        cleaned = re.sub(r'```json\s*', '', response_content)
                        cleaned = re.sub(r'```\s*', '', cleaned)
                        json_start = cleaned.find("{")
                        json_end = cleaned.rfind("}") + 1
                        if json_start >= 0 and json_end > json_start:
                            json_str = cleaned[json_start:json_end]
                            try:
                                result_data = json.loads(json_str)
                            except json.JSONDecodeError:
                                try:
                                    import ast
                                    result_data = ast.literal_eval(json_str)
                                except:
                                    fixed_json = json_str
                                    open_braces = fixed_json.count('{') - fixed_json.count('}')
                                    open_brackets = fixed_json.count('[') - fixed_json.count(']')
                                    fixed_json += ']' * open_brackets + '}' * open_braces
                                    try:
                                        result_data = json.loads(fixed_json)
                                        logger.info(f"  ✅ Fixed truncated JSON by adding {open_brackets} ] and {open_braces} }}")
                                    except:
                                        raise
                            result_map = {r.get('id'): r for r in result_data.get('results', [])}
                            return [{'id': c.get('id'), 'supported': result_map.get(c.get('id'), {}).get('supported', False), 'source_found': True, 'explanation': result_map.get(c.get('id'), {}).get('explanation', '')} for c in claims]
                    except json.JSONDecodeError as e:
                        logger.error(f"  ❌ Batch API JSON parse error: {e}")
                        logger.error(f"  Raw response (first 500 chars): {response_content[:500]}")
                else:
                    logger.warning(f"  ⚠️ Empty response from API for batch verification")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'Parse error: response={response_content[:100] if response_content else "empty"}'} for c in claims]
            except Exception as e:
                if '429' in str(e).lower() or 'rate' in str(e).lower():
                    if retry < API_MAX_RETRIES - 1:
                        time.sleep(API_RETRY_DELAY)
                        continue
                logger.error(f"  ❌ Batch API error: {e}")
                return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': f'API error: {e}'} for c in claims]
        return [{'id': c.get('id'), 'supported': False, 'source_found': True, 'explanation': 'Max retries exceeded'} for c in claims]
    
    def _verify_batch_with_api(self, claims: List[Dict], source_info: Dict) -> List[Dict]:
        if not self.api_client:
            return [{'id': c.get('id'), 'supported': False, 'explanation': 'API not initialized'} for c in claims]
        return self._verify_batch_with_api_thread(claims, source_info, self.api_client)

def main():
    parser = argparse.ArgumentParser(description="Evaluate factual accuracy using LLM agent")
    parser.add_argument("--result", type=str, required=True, help="Path to result markdown file")
    parser.add_argument("--long-context", type=str, help="Path to long context JSON")
    parser.add_argument("--source-folder", type=str, help="Path to source documents folder")
    parser.add_argument("--output", type=str, help="Output file for evaluation result")
    args = parser.parse_args()
    
    result_path = Path(args.result)
    result_text = result_path.read_text(encoding='utf-8')
    evaluator = FactualAccuracyAgentEvaluator()
    result = evaluator.evaluate(result_text=result_text, long_context_path=Path(args.long_context) if args.long_context else None, source_folder=Path(args.source_folder) if args.source_folder else None)
    
    print(f"\nFactual Accuracy Score: {result.score:.1f}/100")
    print(f"Details: {json.dumps(result.details, ensure_ascii=False, indent=2)}")
    
    if args.output:
        output_path = Path(args.output)
        output_path.write_text(result.to_json(), encoding='utf-8')
        print(f"\nResult saved to: {output_path}")

if __name__ == "__main__":
    main()
